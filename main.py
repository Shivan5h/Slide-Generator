from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel
import os
from typing import List, Optional
import uuid
from pptx import Presentation
from pptx.util import Inches
import requests
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
)
import torch
from transformers import pipeline
from huggingface_hub import InferenceClient
from diffusers import StableDiffusionPipeline
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="AI-Powered Slide Generator API",
    description="API for generating PowerPoint presentations using AI with RAG capabilities",
    version="1.0",
)

# Configuration
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Models initialization (using free/open-source options)
EMBEDDING_MODEL = SentenceTransformer("all-MiniLM-L6-v2")
SUMMARIZATION_MODEL = "Falconsai/text_summarization"
LLM_MODEL = "mistralai/Mistral-7B-Instruct-v0.1"
IMAGE_MODEL = "stabilityai/stable-diffusion-xl-base-1.0"

# Initialize components
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
summarizer = pipeline("summarization", model=SUMMARIZATION_MODEL)
llm_client = InferenceClient(model=LLM_MODEL)

# For production, consider using smaller models or implementing caching
try:
    image_pipe = StableDiffusionPipeline.from_pretrained(
        IMAGE_MODEL, torch_dtype=torch.float16
    )
    if torch.cuda.is_available():
        image_pipe = image_pipe.to("cuda")
except Exception as e:
    logger.warning(f"Could not load image generation model: {e}")
    image_pipe = None

class SlideGenerationRequest(BaseModel):
    prompt: str
    document_ids: Optional[List[str]] = None
    generate_images: Optional[bool] = False
    slide_count: Optional[int] = 5
    style: Optional[str] = "professional"

class DocumentUploadResponse(BaseModel):
    document_id: str
    filename: str
    status: str

@app.post("/upload-document/", response_model=DocumentUploadResponse)
async def upload_document(file: UploadFile = File(...)):
    try:
        file_ext = os.path.splitext(file.filename)[1].lower()
        valid_extensions = [".pdf", ".docx", ".txt", ".md"]
        
        if file_ext not in valid_extensions:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        document_id = str(uuid.uuid4())
        file_path = os.path.join(UPLOAD_DIR, f"{document_id}{file_ext}")
        
        with open(file_path, "wb") as f:
            f.write(file.file.read())
        
        return {
            "document_id": document_id,
            "filename": file.filename,
            "status": "uploaded"
        }
    except Exception as e:
        logger.error(f"Document upload failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-slides/")
async def generate_slides(request: SlideGenerationRequest):
    try:
        # Step 1: Process documents if provided
        documents = []
        if request.document_ids:
            for doc_id in request.document_ids:
                doc_path = self._find_document(doc_id)
                if doc_path:
                    loader = self._get_document_loader(doc_path)
                    docs = loader.load()
                    documents.extend(docs)
        
        # Step 2: Create vector store if documents exist
        if documents:
            texts = text_splitter.split_documents(documents)
            text_contents = [t.page_content for t in texts]
            embeddings = EMBEDDING_MODEL.encode(text_contents)
            vector_store = FAISS.from_embeddings(
                list(zip(text_contents, embeddings)), EMBEDDING_MODEL
            )
            
            # Retrieve relevant chunks
            relevant_docs = vector_store.similarity_search(request.prompt, k=3)
            context = "\n\n".join([d.page_content for d in relevant_docs])
        else:
            context = ""
        
        # Step 3: Generate content with LLM
        llm_prompt = f"""
        Create a {request.slide_count}-slide presentation about: {request.prompt}
        {f"Using this context: {context}" if context else ""}
        Return ONLY the presentation outline in this format:
        Slide 1: [Title]\n- [Bullet 1]\n- [Bullet 2]
        Slide 2: [Title]\n- [Bullet 1]\n- [Bullet 2]
        ...
        """
        
        response = llm_client.text_generation(
            llm_prompt,
            max_new_tokens=1500,
            temperature=0.7,
        )
        
        # Step 4: Parse response and create PPT
        ppt = self._create_presentation(response, request.generate_images)
        
        # Save and return
        output_path = f"output_{uuid.uuid4()}.pptx"
        ppt.save(output_path)
        
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="generated_presentation.pptx",
        )
    except Exception as e:
        logger.error(f"Slide generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

def _find_document(self, doc_id: str) -> Optional[str]:
    for f in os.listdir(UPLOAD_DIR):
        if f.startswith(doc_id):
            return os.path.join(UPLOAD_DIR, f)
    return None

def _get_document_loader(self, file_path: str):
    if file_path.endswith(".pdf"):
        return PyPDFLoader(file_path)
    elif file_path.endswith(".docx"):
        return Docx2txtLoader(file_path)
    elif file_path.endswith(".txt"):
        return TextLoader(file_path)
    elif file_path.endswith(".md"):
        return UnstructuredMarkdownLoader(file_path)
    else:
        raise ValueError("Unsupported file type")

def _create_presentation(self, llm_response: str, generate_images: bool) -> Presentation:
    ppt = Presentation()
    
    # Add title slide
    title_slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Generated Presentation"
    subtitle.text = "Created with RAG-powered AI"
    
    # Parse LLM response
    slides_data = []
    current_slide = None
    
    for line in llm_response.split("\n"):
        if line.startswith("Slide"):
            if current_slide:
                slides_data.append(current_slide)
            current_slide = {"title": "", "content": []}
            current_slide["title"] = line.split(":")[1].strip()
        elif line.startswith("-"):
            if current_slide:
                current_slide["content"].append(line[2:])
    
    if current_slide:
        slides_data.append(current_slide)
    
    # Add content slides
    for slide_data in slides_data:
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data["title"]
        content.text = "\n".join(slide_data["content"])
        
        # Generate image if requested
        if generate_images and image_pipe:
            try:
                image_prompt = f"A professional illustration about: {slide_data['title']}"
                image = image_pipe(image_prompt).images[0]
                
                # Save temp image
                img_path = f"temp_{uuid.uuid4()}.png"
                image.save(img_path)
                
                # Add to slide
                left = Inches(6)
                top = Inches(1.5)
                slide.shapes.add_picture(img_path, left, top, height=Inches(3.5))
                
                # Clean up
                os.remove(img_path)
            except Exception as e:
                logger.warning(f"Image generation failed: {e}")
    
    return ppt

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)